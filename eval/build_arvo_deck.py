"""ARVO 2026 poster slide deck for the ENVISION discovery pipeline.

Generates a Google-Slides-importable .pptx with the content Bhavesh
asked for: method schematic, source breakdown, validation metrics,
file-type composition.

Slide order:
  1. Title
  2. Method schematic (the rendered PNG, captioned)
  3. Source repositories table + per-source EYE_IMAGING counts
  4. Validation metrics
  5. File-type composition pie (the rendered PNG, captioned)
  6. Numbers headline + EP integration note

Run:  python eval/build_arvo_deck.py
Out:  paper/envision_arvo_deck.pptx
"""

from __future__ import annotations

import json
import glob
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

ROOT = Path(__file__).resolve().parent.parent
SCHEMATIC = ROOT / "paper" / "envision_method_schematic.png"
PIE = ROOT / "paper" / "envision_filetype_pie.png"
OUT = ROOT / "paper" / "envision_arvo_deck.pptx"

NAVY = RGBColor(0x1D, 0x3A, 0x5E)
TEAL = RGBColor(0x3B, 0x6F, 0xA8)
SLATE = RGBColor(0x55, 0x55, 0x55)
LIGHT = RGBColor(0xEE, 0xF2, 0xF7)


def _per_source_counts() -> list[tuple[str, int, int]]:
    """Return [(source, eye_imaging, total), ...] from latest results."""
    out = []
    for path in sorted(glob.glob(str(ROOT / "results" / "*_all_results.json"))):
        src = Path(path).stem.replace("_all_results", "")
        with open(path) as f:
            recs = json.load(f)
        eye = sum(1 for r in recs if r.get("label") == "EYE_IMAGING")
        out.append((src, eye, len(recs)))
    out.sort(key=lambda t: -t[1])
    return out


def _set_text(tf, text, size=18, bold=False, color=NAVY,
              align=PP_ALIGN.LEFT, italic=False):
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color


def _add_text(slide, text, left, top, width, height, **kw):
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height),
    )
    box.text_frame.word_wrap = True
    _set_text(box.text_frame, text, **kw)
    return box


def _add_bullets(slide, bullets, left, top, width, height,
                 size=14, color=NAVY):
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height),
    )
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = "•  " + b
        run.font.size = Pt(size)
        run.font.color.rgb = color
    return box


def _add_band(slide, color, top_in, height_in=0.85):
    """Coloured horizontal band across the slide."""
    band = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(top_in),
        Inches(13.333), Inches(height_in),
    )
    band.fill.solid()
    band.fill.fore_color.rgb = color
    band.line.fill.background()
    return band


def slide_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _add_band(slide, NAVY, top_in=0.0, height_in=7.5)

    _add_text(
        slide,
        "ENVISION discovery pipeline",
        left=0.5, top=2.0, width=12.3, height=1.5,
        size=54, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
        align=PP_ALIGN.CENTER,
    )
    _add_text(
        slide,
        "Multi-source harvesting · archive-aware metadata · "
        "few-shot classification",
        left=0.5, top=3.4, width=12.3, height=0.8,
        size=22, color=RGBColor(0xCF, 0xDD, 0xED),
        italic=True, align=PP_ALIGN.CENTER,
    )
    _add_text(
        slide,
        "ENVISION Portal  ·  EyeACT",
        left=0.5, top=5.5, width=12.3, height=0.6,
        size=16, color=RGBColor(0xCF, 0xDD, 0xED),
        align=PP_ALIGN.CENTER,
    )
    _add_text(
        slide,
        "github.com/EyeACT/envision-discovery   ·   "
        "github.com/EyeACT/envision-classifier",
        left=0.5, top=6.2, width=12.3, height=0.5,
        size=12, color=RGBColor(0xAA, 0xC0, 0xD7),
        align=PP_ALIGN.CENTER,
    )


def slide_schematic(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_text(
        slide,
        "Method",
        left=0.5, top=0.3, width=12.3, height=0.7,
        size=32, bold=True, color=NAVY,
    )
    if SCHEMATIC.exists():
        slide.shapes.add_picture(
            str(SCHEMATIC),
            Inches(0.4), Inches(1.1),
            width=Inches(12.5),
        )
    _add_text(
        slide,
        "Each scientific dataset record (title + description + keywords + "
        "file types) is encoded with sentence-transformers/all-mpnet-base-v2 "
        "and classified by a SetFit head trained on 891 curated examples "
        "(262 EYE_IMAGING, 629 NEGATIVE).",
        left=0.5, top=6.6, width=12.3, height=0.7,
        size=12, color=SLATE, italic=True, align=PP_ALIGN.CENTER,
    )


def slide_sources(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_text(
        slide,
        "Source repositories",
        left=0.5, top=0.3, width=12.3, height=0.7,
        size=32, bold=True, color=NAVY,
    )
    _add_text(
        slide,
        "Seven public repositories harvested with 47 ophthalmology-specific "
        "queries (OCT, fundus, OCTA, retina, glaucoma, macula, choroid, …)",
        left=0.5, top=1.05, width=12.3, height=0.5,
        size=14, color=SLATE, italic=True,
    )

    rows = _per_source_counts()
    headers = ["Source", "EYE_IMAGING", "Total records", "%"]
    n_rows = len(rows) + 2  # header + body + total

    table_left, table_top = 1.5, 1.7
    table_w, table_h = 10.3, 4.6
    shape = slide.shapes.add_table(
        n_rows, len(headers),
        Inches(table_left), Inches(table_top),
        Inches(table_w), Inches(table_h),
    )
    table = shape.table
    # Column widths (approximate)
    table.columns[0].width = Inches(3.4)
    table.columns[1].width = Inches(2.4)
    table.columns[2].width = Inches(2.4)
    table.columns[3].width = Inches(2.1)

    # Header
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        _set_text(cell.text_frame, h, size=14, bold=True,
                  color=RGBColor(0xFF, 0xFF, 0xFF),
                  align=PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT)

    SOURCE_LABELS = {
        "zenodo": "Zenodo",
        "figshare": "Figshare",
        "datacite": "DataCite",
        "nei": "NIH RePORTER (NEI)",
        "kaggle": "Kaggle",
        "osf": "Open Science Framework",
        "dryad": "Dryad",
    }

    total_eye = sum(e for _, e, _ in rows)
    total_all = sum(t for _, _, t in rows)

    for i, (src, eye, total) in enumerate(rows, start=1):
        pct = (eye / total * 100) if total else 0
        cells = [
            SOURCE_LABELS.get(src, src.title()),
            f"{eye:,}",
            f"{total:,}",
            f"{pct:.1f}%",
        ]
        for j, val in enumerate(cells):
            cell = table.cell(i, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = (
                LIGHT if i % 2 == 1 else RGBColor(0xFF, 0xFF, 0xFF)
            )
            _set_text(cell.text_frame, val, size=14,
                      align=PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT)

    # Total row
    total_pct = (total_eye / total_all * 100) if total_all else 0
    total_cells = ["Total", f"{total_eye:,}", f"{total_all:,}", f"{total_pct:.1f}%"]
    for j, val in enumerate(total_cells):
        cell = table.cell(n_rows - 1, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = TEAL
        _set_text(cell.text_frame, val, size=14, bold=True,
                  color=RGBColor(0xFF, 0xFF, 0xFF),
                  align=PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT)


def slide_metrics(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_text(
        slide,
        "Validation",
        left=0.5, top=0.3, width=12.3, height=0.7,
        size=32, bold=True, color=NAVY,
    )

    metrics = [
        ("0.961", "test accuracy",
         "binary classification on a held-out 20% split"),
        ("0.936", "EYE_IMAGING F1",
         "precision/recall balanced on the positive class"),
        ("90.9%", "manual spot-check",
         "30/33 EYE_IMAGING predictions confirmed by hand"),
        ("891",   "training examples",
         "262 EYE_IMAGING + 629 NEGATIVE, manually curated"),
    ]

    # 2x2 grid of stat cards
    card_w, card_h = 6.0, 2.3
    gap_x, gap_y = 0.5, 0.35
    start_x = (13.333 - 2 * card_w - gap_x) / 2
    start_y = 1.2

    for i, (val, lbl, sub) in enumerate(metrics):
        r, c = divmod(i, 2)
        x = start_x + c * (card_w + gap_x)
        y = start_y + r * (card_h + gap_y)
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(y), Inches(card_w), Inches(card_h),
        )
        card.fill.solid()
        card.fill.fore_color.rgb = LIGHT
        card.line.color.rgb = TEAL
        card.line.width = Pt(1.5)
        # Clear default shape text
        card.text_frame.clear()

        _add_text(
            slide, val,
            left=x, top=y + 0.15, width=card_w, height=0.9,
            size=42, bold=True, color=NAVY, align=PP_ALIGN.CENTER,
        )
        _add_text(
            slide, lbl,
            left=x, top=y + 0.95, width=card_w, height=0.5,
            size=17, bold=True, color=TEAL, align=PP_ALIGN.CENTER,
        )
        _add_text(
            slide, sub,
            left=x, top=y + 1.45, width=card_w, height=0.7,
            size=12, color=SLATE, italic=True, align=PP_ALIGN.CENTER,
        )

    _add_text(
        slide,
        "Ongoing: ophthalmologist Yes/No/I-can't-tell survey on the public "
        "EnvisionPortal validation site (100-record batches; inter-rater "
        "agreement reported once responses accumulate).",
        left=1.0, top=6.7, width=11.3, height=0.6,
        size=12, color=SLATE, italic=True, align=PP_ALIGN.CENTER,
    )


def slide_filetypes(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_text(
        slide,
        "What's actually in the EYE_IMAGING records?",
        left=0.5, top=0.3, width=12.3, height=0.7,
        size=30, bold=True, color=NAVY,
    )
    if PIE.exists():
        slide.shapes.add_picture(
            str(PIE),
            Inches(0.6), Inches(1.1),
            width=Inches(8.5),
        )

    notes = [
        "Allow-list curated for clinical eye-imaging:",
        "JPEG / PNG / TIFF / BMP — fundus, OCT exports",
        "DICOM (.dcm) — medical imaging standard",
        "OCT-native (.e2e Heidelberg, .oct)",
        "Video (.mp4 / .mov / .avi / .cine) — slit-lamp, OCTA",
        "Medical volume (.nii / .mha)",
        "MATLAB (.mat) — common OCT array container",
        "",
        "Excluded: documents, archives, code, microscopy "
        "(.lsm/.lif/.czi/…), animations (.gif/.eps/.svg).",
    ]
    box = slide.shapes.add_textbox(
        Inches(9.3), Inches(1.4), Inches(3.8), Inches(5.5),
    )
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for i, line in enumerate(notes):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.size = Pt(12)
        run.font.color.rgb = NAVY if i == 0 else SLATE
        run.font.bold = (i == 0)
        run.font.italic = (i == len(notes) - 1)


def slide_numbers(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_band(slide, NAVY, top_in=0.0, height_in=7.5)

    _add_text(
        slide, "By the numbers",
        left=0.5, top=0.4, width=12.3, height=0.8,
        size=32, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
        align=PP_ALIGN.CENTER,
    )

    headlines = [
        ("42,524", "records harvested across 7 repositories"),
        ("6,602",  "classified EYE_IMAGING"),
        ("1,896",  "currently registered on EnvisionPortal"),
    ]
    card_w, card_h = 3.8, 4.0
    gap = 0.6
    start_x = (13.333 - 3 * card_w - 2 * gap) / 2
    start_y = 1.8

    for i, (val, lbl) in enumerate(headlines):
        x = start_x + i * (card_w + gap)
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(start_y), Inches(card_w), Inches(card_h),
        )
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        card.line.fill.background()
        card.text_frame.clear()

        _add_text(
            slide, val,
            left=x, top=start_y + 0.6, width=card_w, height=1.6,
            size=72, bold=True, color=NAVY, align=PP_ALIGN.CENTER,
        )
        _add_text(
            slide, lbl,
            left=x, top=start_y + 2.4, width=card_w, height=1.4,
            size=14, color=SLATE, align=PP_ALIGN.CENTER,
        )

    _add_text(
        slide,
        "Difference between 6,602 and 1,896: ENVISION classifies the full "
        "harvested corpus; EnvisionPortal currently shows the records that "
        "have been ingested through its data-conversion pipeline (in progress).",
        left=1.0, top=6.2, width=11.3, height=0.9,
        size=13, color=RGBColor(0xCF, 0xDD, 0xED),
        italic=True, align=PP_ALIGN.CENTER,
    )


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide_title(prs)
    slide_schematic(prs)
    slide_sources(prs)
    slide_metrics(prs)
    slide_filetypes(prs)
    slide_numbers(prs)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"Saved: {OUT}  ({OUT.stat().st_size/1024:.0f} KB)")


if __name__ == "__main__":
    main()
