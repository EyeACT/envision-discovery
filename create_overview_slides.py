#!/usr/bin/env python3
"""Generate ENVISION classifier overview slides for team meeting."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Brand colors
DARK_BG = RGBColor(0x1B, 0x2A, 0x4A)
ACCENT = RGBColor(0x00, 0x96, 0xD6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xE8, 0xEC, 0xF1)
SOFT_GREEN = RGBColor(0x4C, 0xAF, 0x50)
ORANGE = RGBColor(0xFF, 0x98, 0x00)
DARK_TEXT = RGBColor(0x26, 0x32, 0x38)

def add_bg(slide, color=DARK_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_box(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape

def set_text(shape, text, size=14, bold=False, color=WHITE, alignment=PP_ALIGN.LEFT):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    return tf

def add_bullet(tf, text, size=14, color=WHITE, bold=False, level=0, space_before=Pt(6)):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.level = level
    p.space_before = space_before
    return p

# ============================================================
# SLIDE 1: Model Overview
# ============================================================
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide1)

# Title
title = slide1.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11), Inches(0.8))
set_text(title, "ENVISION: Eye Imaging Dataset Classifier", size=32, bold=True, color=WHITE)

# Subtitle
sub = slide1.shapes.add_textbox(Inches(0.8), Inches(1.1), Inches(11), Inches(0.5))
set_text(sub, "Model Overview for Human Expert Evaluation", size=18, color=ACCENT)

# --- What Model ---
box1 = add_box(slide1, Inches(0.8), Inches(2.0), Inches(5.8), Inches(4.8), RGBColor(0x23, 0x3A, 0x5E))
tf1 = set_text(box1, "What Model?", size=22, bold=True, color=ACCENT)
box1.text_frame.paragraphs[0].space_after = Pt(10)

add_bullet(tf1, "SetFit (Sentence Transformer Fine-Tuning)", size=16, bold=True)
add_bullet(tf1, "Few-shot framework from Hugging Face", size=14, level=1, color=LIGHT_GRAY)
add_bullet(tf1, "", size=8)
add_bullet(tf1, "Backbone: sentence-transformers/all-mpnet-base-v2", size=16, bold=True)
add_bullet(tf1, "768-dim embeddings, 384 token context", size=14, level=1, color=LIGHT_GRAY)
add_bullet(tf1, "", size=8)
add_bullet(tf1, "4-Class Output:", size=16, bold=True)
add_bullet(tf1, "EYE_IMAGING  \u2014  actual imaging datasets", size=14, level=1, color=LIGHT_GRAY)
add_bullet(tf1, "EYE_SOFTWARE  \u2014  code, models, tools", size=14, level=1, color=LIGHT_GRAY)
add_bullet(tf1, "EDGE_CASE  \u2014  papers, reviews, borderline", size=14, level=1, color=LIGHT_GRAY)
add_bullet(tf1, "NEGATIVE  \u2014  unrelated domains", size=14, level=1, color=LIGHT_GRAY)

# --- Why This Model ---
box2 = add_box(slide1, Inches(6.9), Inches(2.0), Inches(5.8), Inches(4.8), RGBColor(0x23, 0x3A, 0x5E))
tf2 = set_text(box2, "Why SetFit?", size=22, bold=True, color=ACCENT)
box2.text_frame.paragraphs[0].space_after = Pt(10)

add_bullet(tf2, "Few-shot learning  \u2014  works with small training sets", size=16, bold=True)
add_bullet(tf2, "Only 474 curated examples needed", size=14, level=1, color=LIGHT_GRAY)
add_bullet(tf2, "", size=8)
add_bullet(tf2, "No prompting or generation  \u2014  pure embedding + head", size=16, bold=True)
add_bullet(tf2, "Fast inference, deterministic, no hallucination risk", size=14, level=1, color=LIGHT_GRAY)
add_bullet(tf2, "", size=8)
add_bullet(tf2, "Confidence scores  \u2014  calibrated class probabilities", size=16, bold=True)
add_bullet(tf2, "Supports threshold-based review workflows", size=14, level=1, color=LIGHT_GRAY)
add_bullet(tf2, "", size=8)
add_bullet(tf2, "Input = text metadata only", size=16, bold=True)
add_bullet(tf2, "Title + description + keywords \u2192 classification", size=14, level=1, color=LIGHT_GRAY)
add_bullet(tf2, "No image download or processing required", size=14, level=1, color=LIGHT_GRAY)

# ============================================================
# SLIDE 2: Training & Validation Approach
# ============================================================
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide2)

title2 = slide2.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11), Inches(0.8))
set_text(title2, "Training & Validation Approach", size=32, bold=True, color=WHITE)

# --- Training Data ---
box3 = add_box(slide2, Inches(0.8), Inches(1.6), Inches(5.8), Inches(5.3), RGBColor(0x23, 0x3A, 0x5E))
tf3 = set_text(box3, "Training Data (474 examples)", size=22, bold=True, color=ACCENT)
box3.text_frame.paragraphs[0].space_after = Pt(10)

add_bullet(tf3, "EYE_IMAGING: 77 examples", size=16, bold=True, color=SOFT_GREEN)
add_bullet(tf3, "Known benchmarks (IDRiD, REFUGE, AIROGS, RFMiD, ...)", size=13, level=1, color=LIGHT_GRAY)
add_bullet(tf3, "OCT, OCTA, fundus, corneal, slit-lamp datasets", size=13, level=1, color=LIGHT_GRAY)
add_bullet(tf3, "", size=6)
add_bullet(tf3, "EYE_SOFTWARE: 48 examples", size=16, bold=True, color=ACCENT)
add_bullet(tf3, "GitHub repos, model weights, toolboxes", size=13, level=1, color=LIGHT_GRAY)
add_bullet(tf3, "", size=6)
add_bullet(tf3, "EDGE_CASE: 79 examples", size=16, bold=True, color=ORANGE)
add_bullet(tf3, "Review papers, animal studies, non-imaging eye research", size=13, level=1, color=LIGHT_GRAY)
add_bullet(tf3, "", size=6)
add_bullet(tf3, "NEGATIVE: 270 examples", size=16, bold=True, color=RGBColor(0xEF, 0x53, 0x50))
add_bullet(tf3, "Targeted false-positive patterns:", size=13, level=1, color=LIGHT_GRAY)
add_bullet(tf3, "Cardiovascular OCT (IVOCT), industrial OCT/CT", size=13, level=2, color=LIGHT_GRAY)
add_bullet(tf3, "Microscopy, dermoscopy, brain MRI, taxonomy papers", size=13, level=2, color=LIGHT_GRAY)
add_bullet(tf3, "Robotics (hand-eye calibration), acousto-optics", size=13, level=2, color=LIGHT_GRAY)

# --- Validation ---
box4 = add_box(slide2, Inches(6.9), Inches(1.6), Inches(5.8), Inches(5.3), RGBColor(0x23, 0x3A, 0x5E))
tf4 = set_text(box4, "Validation Approach", size=22, bold=True, color=ACCENT)
box4.text_frame.paragraphs[0].space_after = Pt(10)

add_bullet(tf4, "Zenodo dataset corpus", size=16, bold=True)
add_bullet(tf4, "515 records scraped (resource_type=dataset)", size=14, level=1, color=LIGHT_GRAY)
add_bullet(tf4, "Filtered to 514 with actual data files or links", size=14, level=1, color=LIGHT_GRAY)
add_bullet(tf4, "", size=8)
add_bullet(tf4, "File-type pre-filtering", size=16, bold=True)
add_bullet(tf4, "Only classify records with image/medical/archive files", size=14, level=1, color=LIGHT_GRAY)
add_bullet(tf4, "Exclude genomics-only records (GWAS, RNA-seq)", size=14, level=1, color=LIGHT_GRAY)
add_bullet(tf4, "", size=8)
add_bullet(tf4, "White paper cross-check", size=16, bold=True)
add_bullet(tf4, "Known datasets from ENVISION Portal paper verified", size=14, level=1, color=LIGHT_GRAY)
add_bullet(tf4, "", size=8)
add_bullet(tf4, "Next step: human expert evaluation", size=16, bold=True, color=SOFT_GREEN)
add_bullet(tf4, "Experts review EYE_IMAGING predictions", size=14, level=1, color=LIGHT_GRAY)
add_bullet(tf4, "Focus on medium/low confidence for error analysis", size=14, level=1, color=LIGHT_GRAY)

# ============================================================
# SLIDE 3: Results
# ============================================================
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide3)

title3 = slide3.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11), Inches(0.8))
set_text(title3, "Classification Results", size=32, bold=True, color=WHITE)

# --- Classification Breakdown ---
box5 = add_box(slide3, Inches(0.8), Inches(1.6), Inches(3.7), Inches(5.3), RGBColor(0x23, 0x3A, 0x5E))
tf5 = set_text(box5, "Zenodo Corpus (514 records)", size=20, bold=True, color=ACCENT)
box5.text_frame.paragraphs[0].space_after = Pt(12)

add_bullet(tf5, "EYE_IMAGING:    120", size=18, bold=True, color=SOFT_GREEN)
add_bullet(tf5, "EYE_SOFTWARE:    66", size=18, bold=True, color=ACCENT)
add_bullet(tf5, "EDGE_CASE:          3", size=18, bold=True, color=ORANGE)
add_bullet(tf5, "NEGATIVE:         325", size=18, bold=True, color=RGBColor(0xEF, 0x53, 0x50))
add_bullet(tf5, "", size=10)
add_bullet(tf5, "14 records with external dataset links", size=14, color=LIGHT_GRAY)
add_bullet(tf5, "(GitHub, Kaggle, IEEE DataPort, etc.)", size=12, level=1, color=LIGHT_GRAY)

# --- Confidence Distribution ---
box6 = add_box(slide3, Inches(4.8), Inches(1.6), Inches(3.7), Inches(5.3), RGBColor(0x23, 0x3A, 0x5E))
tf6 = set_text(box6, "Confidence (EYE_IMAGING)", size=20, bold=True, color=ACCENT)
box6.text_frame.paragraphs[0].space_after = Pt(12)

add_bullet(tf6, "High (\u2265 0.95):     117  (97.5%)", size=16, bold=True, color=SOFT_GREEN)
add_bullet(tf6, "Strong candidates, likely correct", size=13, level=1, color=LIGHT_GRAY)
add_bullet(tf6, "", size=8)
add_bullet(tf6, "Medium (0.80\u20130.95):   2", size=16, bold=True, color=ORANGE)
add_bullet(tf6, "Likely eye imaging, worth verifying", size=13, level=1, color=LIGHT_GRAY)
add_bullet(tf6, "", size=8)
add_bullet(tf6, "Lower (< 0.80):         1", size=16, bold=True, color=RGBColor(0xEF, 0x53, 0x50))
add_bullet(tf6, "Manual review recommended", size=13, level=1, color=LIGHT_GRAY)
add_bullet(tf6, "", size=14)
add_bullet(tf6, "\u2192 Model is decisive: nearly all", size=14, color=WHITE)
add_bullet(tf6, "   predictions are high-confidence", size=14, color=WHITE)

# --- Key Takeaways ---
box7 = add_box(slide3, Inches(8.8), Inches(1.6), Inches(3.9), Inches(5.3), RGBColor(0x23, 0x3A, 0x5E))
tf7 = set_text(box7, "Before Expert Evaluation", size=20, bold=True, color=ACCENT)
box7.text_frame.paragraphs[0].space_after = Pt(12)

add_bullet(tf7, "Ready for review:", size=16, bold=True, color=SOFT_GREEN)
add_bullet(tf7, "Standalone classifier API shipped", size=14, level=1, color=LIGHT_GRAY)
add_bullet(tf7, "CLI + Python API available", size=14, level=1, color=LIGHT_GRAY)
add_bullet(tf7, "", size=8)
add_bullet(tf7, "Suggested evaluation plan:", size=16, bold=True, color=WHITE)
add_bullet(tf7, "1. Sample from each class", size=14, level=1, color=LIGHT_GRAY)
add_bullet(tf7, "2. Experts label ground truth", size=14, level=1, color=LIGHT_GRAY)
add_bullet(tf7, "3. Measure precision/recall", size=14, level=1, color=LIGHT_GRAY)
add_bullet(tf7, "4. Focus on EYE_IMAGING vs", size=14, level=1, color=LIGHT_GRAY)
add_bullet(tf7, "   EDGE_CASE boundary", size=14, level=1, color=LIGHT_GRAY)
add_bullet(tf7, "", size=8)
add_bullet(tf7, "Open questions:", size=16, bold=True, color=ORANGE)
add_bullet(tf7, "How many records to review?", size=14, level=1, color=LIGHT_GRAY)
add_bullet(tf7, "How many expert reviewers?", size=14, level=1, color=LIGHT_GRAY)
add_bullet(tf7, "Review tool / annotation format?", size=14, level=1, color=LIGHT_GRAY)

# Save
out_path = "/home/joneill/Nextcloud/vaults/jmind/calmi2/envision-discovery/ENVISION_Classifier_Overview.pptx"
prs.save(out_path)
print(f"Saved to: {out_path}")
